import sys
sys.path.insert(0, "/private/tmp/claude/pptx-deps")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT    = RGBColor(0x4A, 0x9E, 0xFF)   # blue
ACCENT2   = RGBColor(0x3D, 0xD6, 0xB5)   # teal
WARN      = RGBColor(0xFF, 0xB8, 0x4C)   # amber
SUCCESS   = RGBColor(0x4C, 0xD9, 0x7A)   # green
DANGER    = RGBColor(0xFF, 0x5C, 0x5C)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DIM       = RGBColor(0x8A, 0x9B, 0xB8)
CODE_BG   = RGBColor(0x1A, 0x24, 0x3A)
CODE_FG   = RGBColor(0xA8, 0xD8, 0xFF)
PURPLE    = RGBColor(0xC0, 0x8A, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # truly blank

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, bg=None, border=None):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if border is None else None
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(sl, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Inter" if not False else "Courier New"
    return txb

def code_box(sl, code_lines, x, y, w, h, size=11):
    shape = box(sl, x, y, w, h, bg=CODE_BG, border=ACCENT)
    txb = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.12), Inches(w-0.3), Inches(h-0.24))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.name = "Courier New"
        run.font.color.rgb = CODE_FG
    return txb

def pill(sl, text, x, y, color=ACCENT, text_color=WHITE, size=11):
    w = len(text) * 0.085 + 0.25
    h = 0.32
    shape = box(sl, x, y, w, h, bg=color)
    shape.adjustments[0] = 0.5
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Inter"

def divider(sl, y, color=ACCENT):
    ln = sl.shapes.add_connector(1, Inches(0.5), Inches(y), Inches(12.83), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(0.75)

def header(sl, title, subtitle=None):
    box(sl, 0, 0, 13.33, 1.1, bg=RGBColor(0x16, 0x22, 0x3A))
    txt(sl, title, 0.45, 0.18, 10, 0.55, size=26, bold=True, color=ACCENT)
    if subtitle:
        txt(sl, subtitle, 0.45, 0.7, 10, 0.38, size=13, color=DIM)

def arrow_right(sl, x, y, length=0.6, color=DIM):
    x1, y1 = Inches(x), Inches(y)
    x2, y2 = Inches(x + length), Inches(y)
    ln = sl.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(1.5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
# big accent bar
box(sl, 0, 0, 0.35, 7.5, bg=ACCENT)
# geometric deco boxes
box(sl, 11.2, 5.8, 2.13, 1.7, bg=RGBColor(0x16, 0x22, 0x3A))
box(sl, 11.8, 0,   1.53, 2.2, bg=RGBColor(0x16, 0x22, 0x3A))

txt(sl, "AGENTIC\nENGINEERING\nPIPELINE", 0.9, 1.2, 9, 3.6,
    size=54, bold=True, color=WHITE)
txt(sl, "How the code works", 0.9, 4.7, 8, 0.7, size=22, color=ACCENT2)
txt(sl, "state machine  ·  orchestrator  ·  agent runner  ·  ledger  ·  retry loop",
    0.9, 5.35, 10, 0.5, size=13, color=DIM)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Architecture Overview", "Three layers — each with one job")

layers = [
    ("state-machine.mjs",  "WHAT to run",   "Declares the 10 states, their agents, success/fail transitions, and human gates. Pure data — no logic.",  ACCENT),
    ("orchestrator.mjs",   "HOW to run it", "Walks the state machine in a while loop. Manages the ledger, retries, gates, and escalation.",             ACCENT2),
    ("claude-cli.mjs",     "RUNS each agent","Spawns  claude -p  as a child process. Assembles the system prompt, scopes tools, parses output.",        PURPLE),
]

for i, (fname, role, desc, col) in enumerate(layers):
    bx = 0.45 + i * 4.22
    box(sl, bx, 1.35, 3.85, 4.9, bg=RGBColor(0x16, 0x22, 0x3A), border=col)
    txt(sl, fname, bx+0.18, 1.55, 3.5, 0.42, size=13, bold=True, color=col)
    txt(sl, role,  bx+0.18, 2.05, 3.5, 0.38, size=16, bold=True, color=WHITE)
    txt(sl, desc,  bx+0.18, 2.55, 3.5, 3.2,  size=11.5, color=DIM, wrap=True)
    if i < 2:
        arrow_right(sl, bx+3.85, 3.8, 0.37, color=ACCENT)

txt(sl, "Intelligence lives in the agents.  Reliability lives in the orchestrator.",
    0.45, 6.6, 12.43, 0.5, size=13, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — State Machine
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "state-machine.mjs", "Declares the pipeline as a graph — no logic, pure data")

code_box(sl, [
    "export const STATES = {",
    "  scout:         { agent: 'scout',         gate: null,                 next: 'spec'          },",
    "  spec:          { agent: 'spec-writer',   gate: 'human-approve-spec', next: 'plan'          },",
    "  plan:          { agent: 'planner',       gate: null,                 next: 'implement'     },",
    "  implement:     { agent: 'implementer',   gate: null,                 next: 'test',    onFail: 'implement', isolation: 'worktree' },",
    "  test:          { agent: 'test-author',   gate: null,                 next: 'pr',      onFail: 'implement' },",
    "  pr:            { agent: 'pr-agent',      gate: null,                 next: 'review'        },",
    "  review:        { agent: 'reviewer',      gate: null,                 next: 'preview_e2e',  onFail: 'implement' },",
    "  preview_e2e:   { agent: 'preview-e2e',   gate: 'human-approve-pr',   next: 'merge_release' },",
    "  merge_release: { agent: 'merge-release', gate: null,                 next: 'close_curate', onFail: 'ROLLBACK' },",
    "  close_curate:  { agent: 'closer-curator',gate: null,                 next: 'DONE'          },",
    "}",
], 0.45, 1.25, 12.43, 4.55, size=10)

# callout boxes
items = [
    (0.45,  6.1, "agent",     "which function runs",   ACCENT),
    (3.2,   6.1, "next",      "success → go here",     ACCENT2),
    (5.95,  6.1, "onFail",    "failure → route back",  WARN),
    (8.7,   6.1, "gate",      "pause for human",       PURPLE),
    (11.2,  6.1, "isolation", "own git worktree",       SUCCESS),
]
for x, y, label, desc, col in items:
    box(sl, x, y, 2.45, 0.9, bg=RGBColor(0x1A, 0x24, 0x3A), border=col)
    txt(sl, label, x+0.12, y+0.08, 2.2, 0.32, size=11, bold=True, color=col)
    txt(sl, desc,  x+0.12, y+0.44, 2.2, 0.38, size=10, color=DIM)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Ledger
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "The Ledger", "The single source of truth — one object, shared across every agent")

code_box(sl, [
    "// Created once in orchestrator.mjs before the loop starts",
    "const ledger = {",
    "  ticketKey: 'TASK-142',",
    "  request:   'We want a deadline-based task scheduler...',  // original feature request",
    "  state:     'scout',",
    "  attempts:  {},        // { scout: 1, implement: 2, ... }",
    "  artifacts: {},        // populated by each agent as it completes",
    "  history:   [],        // audit trail of every agent run",
    "}",
    "",
    "// After the pipeline runs, artifacts looks like:",
    "ledger.artifacts = {",
    "  scout:         '# Impact assessment...',                   // markdown string",
    "  spec:          { title, problem, acceptanceCriteria },     // parsed JSON",
    "  plan:          { slices: [{ backend }, { frontend }] },    // parsed JSON",
    "  implement:     '# Implementation summary...',             // markdown string",
    "  test:          '# Test report...',                         // markdown string",
    "  pr:            { description, diff },                      // parsed JSON",
    "  review:        { verdict: 'pass', findings: [...] },       // parsed JSON",
    "}",
], 0.45, 1.25, 12.43, 5.6, size=10.5)

txt(sl, "No database. No message queue. Just this object in memory — passed by reference to every agent function.",
    0.45, 7.0, 12.43, 0.38, size=12, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Orchestrator Loop
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "orchestrator.mjs", "A while loop that walks the state machine")

code_box(sl, [
    "while (ledger.state !== 'DONE' && ledger.state !== 'ESCALATE') {",
    "  const stateId = ledger.state;                    // e.g. 'scout'",
    "  const def     = STATES[stateId];                 // { agent, next, onFail, gate }",
    "  const agentFn = agents[def.agent];               // the actual function to call",
    "",
    "  ledger.attempts[stateId] = (ledger.attempts[stateId] ?? 0) + 1;",
    "",
    "  // ── Run the agent ──────────────────────────────────────────────────",
    "  const result = await agentFn({ ledger, attempt, log });",
    "  //                              ^^^^^^",
    "  //                              whole ledger passed in — agent reads what it needs",
    "",
    "  if (result.artifact) ledger.artifacts[stateId] = result.artifact;",
    "  //                                    ^^^^^^^   one entry added per completed state",
    "",
    "  // ── Route on failure ──────────────────────────────────────────────",
    "  if (result.ok === false) {",
    "    if (attempt >= MAX_ATTEMPTS) { ledger.state = 'ESCALATE'; break; }",
    "    ledger.state = failState(stateId);   // e.g. 'test' → back to 'implement'",
    "    continue;",
    "  }",
    "",
    "  // ── Human gate (blocks until approved) ───────────────────────────",
    "  if (def.gate) {",
    "    const decision = await approveGate(def.gate, ledger);",
    "    if (!decision.approved) { ledger.state = 'ESCALATE'; break; }",
    "  }",
    "",
    "  ledger.state = nextState(stateId);     // advance to next state",
    "}",
], 0.45, 1.25, 12.43, 5.95, size=9.8)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — How an Agent Runs (runClaude)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "claude-cli.mjs  —  runClaude()", "Spawns  claude -p  as a child process for every agent turn")

code_box(sl, [
    "export async function runClaude({ prompt, agentPromptPath, allowedTools, schema, permissionMode }) {",
    "",
    "  // 1. Build system prompt: shared knowledge base + this agent's role file",
    "  const systemPrompt = KNOWLEDGE + '\\n\\n---\\n\\n' + readFileSync(agentPromptPath);",
    "  //    KNOWLEDGE = .knowledge/CLAUDE.md + skills/*.md + decisions/*.md",
    "  //    agentPromptPath = e.g. 'agents/07-reviewer.md'",
    "",
    "  // 2. If structured output needed, append JSON Schema to the user prompt",
    "  const finalPrompt = schema",
    "    ? `${prompt}\\n\\nRespond with ONLY a JSON object matching:\\n${JSON.stringify(schema)}`",
    "    : prompt;",
    "",
    "  // 3. Spawn the local Claude Code CLI in headless mode",
    "  const { stdout } = await execFile('claude', [",
    "    '-p',                   finalPrompt,",
    "    '--output-format',      'json',",
    "    '--append-system-prompt', systemPrompt,",
    "    '--permission-mode',    permissionMode,   // 'plan' | 'acceptEdits'",
    "    '--allowedTools',       allowedTools.join(','),",
    "  ]);",
    "",
    "  // 4. Parse result — raw text or JSON object",
    "  const env = JSON.parse(stdout);",
    "  return schema ? coerceJson(env.result) : env.result;",
    "}",
], 0.45, 1.25, 12.43, 5.8, size=10)

# bottom callout row
callouts = [
    (0.45,  7.05, "No API key",     "Uses your local  claude  login",      ACCENT2),
    (4.0,   7.05, "Stateless",      "Each call is a fresh subprocess",      ACCENT),
    (7.55,  7.05, "Scoped tools",   "Per-agent allowlist",                  WARN),
    (11.1,  7.05, "Structured out", "JSON Schema forces typed return",      PURPLE),
]
for x, y, label, desc, col in callouts:
    txt(sl, f"▸ {label}", x, y,    3.3, 0.28, size=11, bold=True, color=col)
    txt(sl, desc,          x, y+0.28, 3.3, 0.22, size=10, color=DIM)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — System Prompt Assembly
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "System Prompt Assembly", "Every agent gets the same knowledge base + its own role on top")

# diagram
layers_data = [
    (".knowledge/CLAUDE.md",         "Coding conventions, UTC rules, naming standards",    ACCENT,  1.45),
    (".knowledge/skills/*.md",        "How to add an Angular feature, NestJS module, etc.", ACCENT,  2.3),
    (".knowledge/decisions/*.md",     "ADRs — e.g. ADR-0001: scheduler uses cron",          ACCENT,  3.15),
    ("─── separator ─────────────",  "",                                                    DIM,     3.82),
    ("agents/07-reviewer.md",         "Role: reviewer. Tools allowed. Output format.",       PURPLE,  4.25),
]
for label, desc, col, y in layers_data:
    if label.startswith("───"):
        divider(sl, y, color=DIM)
        txt(sl, "appended via  --append-system-prompt", 5.5, y-0.2, 4, 0.35, size=10, color=DIM, align=PP_ALIGN.CENTER)
        continue
    box(sl, 0.45, y, 7.5, 0.7, bg=RGBColor(0x16, 0x22, 0x3A), border=col)
    txt(sl, label, 0.65, y+0.1,  7.1, 0.3, size=11, bold=True, color=col)
    txt(sl, desc,  0.65, y+0.38, 7.1, 0.26, size=10, color=DIM)

txt(sl, "KNOWLEDGE string", 4.0, 1.25, 4, 0.35, size=13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# right side: code
code_box(sl, [
    "// Loaded ONCE, cached across all agents",
    "function loadKnowledge() {",
    "  const parts = [",
    "    readFileSync('.knowledge/CLAUDE.md'),",
    "  ];",
    "  for (const dir of ['skills', 'decisions']) {",
    "    for (const f of readdirSync(dir))",
    "      parts.push(readFileSync(f));",
    "  }",
    "  return parts.join('\\n\\n---\\n\\n');",
    "}",
    "",
    "// Per-agent prompt built on demand",
    "function systemPrompt(agentPromptPath) {",
    "  const role = readFileSync(agentPromptPath);",
    "  return `${KNOWLEDGE}\\n\\n---\\n\\n${role}`;",
    "}",
], 8.3, 1.45, 4.58, 5.5, size=9.5)

txt(sl, "Claude Code caches the large stable prefix automatically — only the agent role tip changes per call.",
    0.45, 7.05, 12.43, 0.38, size=11.5, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Information Passing
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Information Passing Between Agents", "No shared memory — prior outputs are copy-pasted into the next prompt")

# left: code
code_box(sl, [
    "// SCOUT — reads ledger.request, returns markdown",
    "const scout = async ({ ledger }) => {",
    "  const md = await runClaude({",
    "    prompt: `Feature request:\\n\\n${ledger.request}`,",
    "  });",
    "  return { ok: true, artifact: md };",
    "  // → stored as: ledger.artifacts['scout'] = md",
    "};",
    "",
    "// SPEC-WRITER — reads .request AND .artifacts.scout",
    "const specWriter = async ({ ledger }) => {",
    "  const spec = await runClaude({",
    "    prompt: `Request:\\n${ledger.request}",
    "             Impact:\\n${ledger.artifacts.scout}`,",
    "    //                  ^^^^^^^^^^^^^^^^^^^^",
    "    //                  scout's markdown pasted inline",
    "    schema: SPEC_SCHEMA,",
    "  });",
    "  return { ok: true, artifact: spec };",
    "  // → stored as: ledger.artifacts['spec'] = { title, ACs, ... }",
    "};",
    "",
    "// REVIEWER — reads spec + implement + test + pr at once",
    "const reviewer = async ({ ledger }) => {",
    "  const prompt = `Spec:\\n${JSON.stringify(ledger.artifacts.spec)}",
    "    Impl:\\n${ledger.artifacts.implement}",
    "    Tests:\\n${ledger.artifacts.test}",
    "    Diff:\\n${ledger.artifacts.pr?.diff}`;",
    "};",
], 0.45, 1.25, 6.8, 5.95, size=9.5)

# right: flow diagram
txt(sl, "Data flow", 8.0, 1.3, 5, 0.38, size=14, bold=True, color=WHITE)

flow = [
    ("ledger.request",          ACCENT,  1.85),
    ("ledger.artifacts.scout",  ACCENT2, 2.55),
    ("ledger.artifacts.spec",   ACCENT2, 3.25),
    ("ledger.artifacts.plan",   ACCENT2, 3.95),
    ("ledger.artifacts.implement", WARN, 4.65),
    ("ledger.artifacts.test",   WARN,    5.35),
    ("ledger.artifacts.pr",     WARN,    6.05),
]
for label, col, y in flow:
    box(sl, 7.65, y, 5.15, 0.55, bg=RGBColor(0x16, 0x24, 0x3A), border=col)
    txt(sl, label, 7.8, y+0.12, 4.8, 0.32, size=10.5, bold=True, color=col)
    if y < 6.05:
        arrow_right(sl, 10.22, y+0.55, 0, color=DIM)

txt(sl, "The mechanism is string concatenation.\nEach agent reads what it needs from ledger.artifacts\nand pastes it into its prompt as plain text.",
    7.65, 6.75, 5.2, 0.7, size=10.5, italic=True, color=DIM, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Retry Loop
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "The Retry Loop", "ok: false routes back — the orchestrator, not the agent, controls flow")

# left code
code_box(sl, [
    "// orchestrator.mjs — failure routing",
    "if (result.ok === false) {",
    "  const target = failState(stateId);",
    "  // 'test'   → onFail → 'implement'",
    "  // 'review' → onFail → 'implement'",
    "",
    "  if (attempt >= MAX_ATTEMPTS) {   // MAX = 3",
    "    ledger.state = 'ESCALATE';",
    "    break;",
    "  }",
    "  ledger.state = target;",
    "  continue;  // restarts the while loop",
    "}",
    "",
    "// agents.cli.mjs — implementer on retry",
    "const liveImplementer = async ({ ledger, attempt }) => {",
    "  const fix = ledger.testFail",
    "    ? `Tests are FAILING. Fix this:\\n${ledger.testFail.slice(-3500)}`",
    "    : '';",
    "  await runClaude({",
    "    prompt: `Plan:\\n${JSON.stringify(ledger.artifacts.plan)}${fix}`,",
    "    //               ^^^^ failure output appended on retry",
    "  });",
    "};",
], 0.45, 1.25, 6.8, 5.6, size=9.8)

# right: flow diagram
txt(sl, "Retry flow", 8.0, 1.3, 5, 0.38, size=14, bold=True, color=WHITE)

steps = [
    ("implement",     ACCENT,  1.8,  "attempt 1, 2, 3"),
    ("test-author",   ACCENT2, 2.75, "runs npm test"),
    ("ok: false",     DANGER,  3.7,  "CI red → routes back"),
    ("implement",     WARN,    4.65, "attempt 2 — gets failure output"),
    ("test-author",   SUCCESS, 5.6,  "ok: true → continues"),
]
for label, col, y, note in steps:
    box(sl, 8.0, y, 3.2, 0.6, bg=RGBColor(0x16, 0x24, 0x3A), border=col)
    txt(sl, label, 8.15, y+0.12, 3.0, 0.32, size=12, bold=True, color=col)
    txt(sl, note,  11.35, y+0.15, 1.85, 0.32, size=9.5, color=DIM)
    if y < 5.6:
        arrow_right(sl, 9.6, y+0.6, 0, color=DIM)

txt(sl, f"After MAX_ATTEMPTS (3) the orchestrator sets\nledger.state = 'ESCALATE' and stops the loop.",
    8.0, 6.5, 5.0, 0.65, size=10.5, italic=True, color=DIM, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Human Gates
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Human Gates", "The loop pauses — nothing moves forward until a person approves")

code_box(sl, [
    "// orchestrator.mjs — gate check runs AFTER a successful agent result",
    "if (def.gate) {",
    "  log.gate(def.gate);                          // '✋ HUMAN GATE: human-approve-spec'",
    "  const decision = await approveGate(def.gate, ledger);",
    "  //               ^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^",
    "  //               In production: blocks on a real Linear approval or GitHub PR review.",
    "  //               In the demo:   auto-approves immediately.",
    "",
    "  if (!decision.approved) {",
    "    ledger.state = 'ESCALATE';                 // human rejected → stop pipeline",
    "    break;",
    "  }",
    "}",
    "// Only reaches nextState() if the gate cleared",
    "ledger.state = nextState(stateId);",
], 0.45, 1.25, 12.43, 3.6, size=11)

# two gates
gates = [
    ("human-approve-spec", "After spec-writer", "A PM reads the ticket and acceptance criteria before any code is written.\nThe spec + open questions are in ledger.artifacts.spec.", ACCENT),
    ("human-approve-pr",   "After preview-e2e", "A tech lead reviews the PR and E2E report before merge.\nAll artifacts are available in ledger.artifacts at this point.", PURPLE),
]
for i, (name, when, desc, col) in enumerate(gates):
    x = 0.45 + i * 6.44
    box(sl, x, 5.05, 6.0, 2.1, bg=RGBColor(0x16, 0x22, 0x3A), border=col)
    txt(sl, name, x+0.2, 5.15, 5.6, 0.4, size=13, bold=True, color=col)
    txt(sl, when, x+0.2, 5.55, 5.6, 0.3, size=11, color=DIM)
    txt(sl, desc, x+0.2, 5.95, 5.6, 1.0, size=10.5, color=WHITE, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Permission Model
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Permission Model", "Least-privilege — each agent only gets the tools its role needs")

rows = [
    ("scout",          "plan      (read-only)",  "Read, Grep, Glob",                     "Explores codebase — never writes",          ACCENT),
    ("spec-writer",    "plan      (read-only)",  "Read, mcp.linear.write",               "Reads repo, creates Linear ticket",         ACCENT),
    ("planner",        "plan      (read-only)",  "Read",                                 "Reads spec, outputs JSON plan",             ACCENT),
    ("implementer",    "acceptEdits",            "Read, Edit, Write, Bash, Grep, Glob",  "The ONLY agent that writes code",           WARN),
    ("test-author",    "acceptEdits",            "Read, Edit, Write, Bash, Grep, Glob",  "Writes tests, runs npm test",               WARN),
    ("pr-agent",       "acceptEdits",            "Read, Bash",                           "Reads diff, writes PR description",         WARN),
    ("reviewer",       "plan      (read-only)",  "Read",                                 "Reads artifacts, emits verdict JSON",       ACCENT2),
    ("closer-curator", "plan      (read-only)",  "Read, mcp.linear.write",               "Closes ticket, proposes CLAUDE.md update",  ACCENT2),
]

hdrs = ["Agent", "permissionMode", "allowedTools", "What it can do"]
col_x = [0.45, 2.6, 4.9, 8.6]
col_w = [2.05, 2.2, 3.6, 4.1]

for j, h in enumerate(hdrs):
    txt(sl, h, col_x[j], 1.3, col_w[j], 0.35, size=11, bold=True, color=DIM)
divider(sl, 1.65, color=DIM)

for i, (agent, pmode, tools, desc, col) in enumerate(rows):
    y = 1.75 + i * 0.62
    if i % 2 == 0:
        box(sl, 0.35, y-0.04, 12.63, 0.58, bg=RGBColor(0x14, 0x1E, 0x32))
    txt(sl, agent, col_x[0], y,      col_w[0], 0.4, size=10.5, bold=True, color=col)
    pcolor = DANGER if "acceptEdits" in pmode else ACCENT2
    txt(sl, pmode, col_x[1], y,      col_w[1], 0.4, size=9.5, color=pcolor)
    txt(sl, tools, col_x[2], y,      col_w[2], 0.4, size=9,   color=DIM)
    txt(sl, desc,  col_x[3], y,      col_w[3], 0.4, size=9.5, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Live vs Simulated
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Live vs Simulated Mode", "Same orchestrator, same state machine — only the agent functions swap")

code_box(sl, [
    "// run.mjs — one env var switches the entire agent layer",
    "const LIVE = process.env.AGENTS === 'live';",
    "const { makeAgents } = await import(",
    "  LIVE ? './agents.cli.mjs'       // real claude -p calls + real npm test",
    "       : './simulated-agents.mjs' // hard-coded artifacts, runs offline",
    ");",
], 0.45, 1.3, 12.43, 2.0, size=12)

cols_data = [
    ("Simulated\nnode example/run.mjs",
     ["Pre-written artifacts returned immediately",
      "One CI failure injected to demo retry routing",
      "Human gates auto-approve",
      "Runs offline — no Claude login needed",
      "Good for understanding the control flow"],
     ACCENT2, 0.45),
    ("Live\nAGENTS=live node example/run.mjs",
     ["scout / spec / plan / review / curate → real claude -p calls",
      "implementer / test-author → edit real files under apps/taskapp",
      "test-author → runs real npm test as the CI gate",
      "pr-agent → generates PR description from real git diff",
      "Linear ticket created/closed if LINEAR_API_KEY set"],
     WARN, 6.7),
]
for label, points, col, rx in cols_data:
    x = 0.45 if rx < 5 else 6.94
    box(sl, x, 3.5, 6.0, 3.65, bg=RGBColor(0x16, 0x22, 0x3A), border=col)
    txt(sl, label, x+0.2, 3.62, 5.6, 0.6, size=12, bold=True, color=col, wrap=True)
    for j, pt in enumerate(points):
        txt(sl, f"▸  {pt}", x+0.2, 4.35+j*0.55, 5.6, 0.48, size=10.5, color=WHITE, wrap=True)

txt(sl, "The orchestrator doesn't know or care which mode is active.",
    0.45, 7.1, 12.43, 0.35, size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Summary
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
box(sl, 0, 0, 0.35, 7.5, bg=ACCENT)
txt(sl, "Summary", 0.9, 0.6, 10, 0.7, size=36, bold=True, color=WHITE)
divider(sl, 1.5, color=ACCENT)

points = [
    (ACCENT,  "state-machine.mjs",   "Pure data. Declares 10 states: which agent runs, where to go on success, where to retry on failure, when to pause for a human."),
    (ACCENT2, "orchestrator.mjs",    "A while loop. Runs each agent, writes its output to ledger.artifacts[stateId], routes back on failure (ok: false), pauses at gates, escalates after MAX_ATTEMPTS."),
    (PURPLE,  "claude-cli.mjs",      "Spawns  claude -p  as a child process. Assembles system prompt = knowledge base + agent role. Scopes tools. Returns raw text or parsed JSON."),
    (WARN,    "The ledger",          "One plain JS object. Created once, passed by reference to every agent. Agents read prior artifacts from it and return their own as { artifact }."),
    (SUCCESS, "Information passing", "String concatenation. Each agent function manually pulls what it needs from ledger.artifacts and pastes it into the prompt string before calling runClaude()."),
    (DANGER,  "Retry loop",          "Failure stays inside the orchestrator. The agent returns ok: false, the orchestrator routes back to the onFail state and appends the error to the next prompt."),
]
for i, (col, label, desc) in enumerate(points):
    y = 1.7 + i * 0.92
    box(sl, 0.75, y, 0.08, 0.55, bg=col)
    txt(sl, label, 1.05, y,      3.2, 0.38, size=12, bold=True, color=col)
    txt(sl, desc,  1.05, y+0.4, 11.6, 0.45, size=10.5, color=DIM, wrap=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "agentic-pipeline-architecture.pptx"
prs.save(out)
print(f"Saved: {out}")
